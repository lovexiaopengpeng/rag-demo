import os
import pdfplumber
from langchain_core.documents import Document
from langchain_chroma import Chroma
from langchain_text_splitters import RecursiveCharacterTextSplitter

# 尝试导入OllamaEmbeddings
try:
    from langchain_ollama import OllamaEmbeddings
except ImportError:
    print("⚠️ 未安装langchain_ollama模块，将使用默认嵌入方法")
    # 使用一个简单的占位符
    class OllamaEmbeddings:
        def __init__(self, model=None):
            pass
        def embed_documents(self, texts):
            return [[0.0]*128 for _ in texts]
        def embed_query(self, text):
            return [0.0]*128

# 尝试导入docx模块
try:
    from docx import Document as DocxDocument
    has_docx = True
except ImportError:
    has_docx = False
    print("⚠️ 未安装python-docx模块，将跳过.docx文件处理")

from pathlib import Path
from typing import List

from rag_core import VECTOR_DB

from rag_core import EMBEDDINGS

import uuid
from datetime import datetime
from typing import List

# 尝试导入图片处理相关库
try:
    from PIL import Image
    import pytesseract
    import cv2
    import numpy as np
    has_image_processing = True
    print("✅ 成功导入图片处理库")
except ImportError as e:
    has_image_processing = False
    print(f"⚠️ 未安装图片处理库，将跳过图片文件处理: {e}")

DOCS_DIR = Path("./docs")#"./docs"
#DB_DIR = "./chroma_db"

embedding = EMBEDDINGS # OllamaEmbeddings(model="nomic-embed-text")

#vectordb = Chroma(
#    persist_directory=str(DB_DIR),
#    embedding_function=embedding
#)

splitter = RecursiveCharacterTextSplitter(
    chunk_size=500,
    chunk_overlap=100
)

#def load_txt(path: str):
#    with open(path, "r", encoding="utf-8") as f:
#        text = f.read()
#    return [Document(page_content=text, metadata={"source": path})]

def load_txt_old(path: str):
    print(f"读取txt文本文件: {path}")
    with open(path, "r", encoding="utf-8") as f:
        text = f.read()
        return [Document(page_content=text, 
                         metadata={"source": os.path.basename(path), "type": "txt"})]
    #return [Document(page_content=text, metadata={"source": path})]

def load_txt(path: Path) -> List[Document]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    return [
        Document(
            page_content=text,
            metadata={"source": path.name,"type": "txt"}
        )
    ]



def load_md(path: str):
    return load_txt(path)

def load_pdf_old(path: str):
    texts = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            if page.extract_text():
                texts.append(page.extract_text())

    content = "\n".join(texts)
    if not content.strip():
        return []

    return [Document(page_content=content, metadata={"source": path.name,"type": "pdf"})]

def load_pdf(path: Path) -> list[Document]:
    docs = []

    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text()
            if not text or not text.strip():
                continue

            docs.append(
                Document(
                    page_content=text,
                    metadata={
                        "source": path.name,  # ✅ 一定是 str
                        "type": "pdf",
                        "page": i + 1         # ✅ int
                    }
                )
            )

    return docs

def load_docx(path: str):
    if not has_docx:
        print("⚠️ 未安装python-docx模块，无法处理.docx文件")
        return []
    doc = DocxDocument(path)
    text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    return [Document(page_content=text, metadata={"source": path, "type": "docx"})]

def load_xlsx(path: Path) -> List[Document]:
    """
    加载Excel文件，提取所有工作表中的文本内容
    """
    from openpyxl import load_workbook
    
    wb = load_workbook(filename=path, read_only=True)
    docs = []
    
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_text = []
        
        # 提取工作表名称
        sheet_text.append(f"工作表: {sheet_name}")
        sheet_text.append("=" * 30)
        
        # 提取所有单元格内容
        for row in ws.iter_rows(values_only=True):
            # 过滤空行
            row_text = [str(cell) if cell is not None else "" for cell in row]
            if any(row_text):
                sheet_text.append("\t".join(row_text))
        
        # 添加工作表内容到文档列表
        full_text = "\n".join(sheet_text)
        if full_text.strip():
            docs.append(
                Document(
                    page_content=full_text,
                    metadata={
                        "source": path.name,
                        "type": "xlsx",
                        "sheet": sheet_name
                    }
                )
            )
    
    return docs

def load_pptx(path: Path) -> List[Document]:
    """
    加载PowerPoint文件，提取所有幻灯片中的文本内容
    """
    from pptx import Presentation
    
    prs = Presentation(path)
    docs = []
    
    for i, slide in enumerate(prs.slides, 1):
        slide_text = []
        
        # 提取幻灯片编号
        slide_text.append(f"幻灯片 {i}")
        slide_text.append("=" * 30)
        
        # 提取幻灯片中的所有文本
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text:
                slide_text.append(shape.text_frame.text)
            elif hasattr(shape, "table"):
                # 处理表格内容
                table = shape.table
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    if any(row_text):
                        slide_text.append("\t".join(row_text))
        
        # 添加幻灯片内容到文档列表
        full_text = "\n".join(slide_text)
        if full_text.strip():
            docs.append(
                Document(
                    page_content=full_text,
                    metadata={
                        "source": path.name,
                        "type": "pptx",
                        "slide": i
                    }
                )
            )
    
    return docs


def load_csv(path: Path) -> List[Document]:
    """
    加载CSV文件，提取文本内容
    """
    import csv
    
    csv_text = []
    csv_text.append(f"CSV文件: {path.name}")
    csv_text.append("=" * 30)
    
    with open(path, 'r', encoding='utf-8', errors='ignore') as f:
        reader = csv.reader(f)
        for row in reader:
            if row:
                csv_text.append("\t".join(row))
    
    full_text = "\n".join(csv_text)
    return [
        Document(
            page_content=full_text,
            metadata={"source": path.name, "type": "csv"}
        )
    ]


def load_html(path: Path) -> List[Document]:
    """
    加载HTML文件，提取文本内容
    """
    from bs4 import BeautifulSoup
    
    html_content = path.read_text(encoding='utf-8', errors='ignore')
    soup = BeautifulSoup(html_content, 'html.parser')
    
    # 提取文本内容，去除多余空白
    text = soup.get_text(separator='\n', strip=True)
    
    return [
        Document(
            page_content=text,
            metadata={"source": path.name, "type": "html"}
        )
    ]


def load_json(path: Path) -> List[Document]:
    """
    加载JSON文件，提取文本内容
    """
    import json
    
    json_content = path.read_text(encoding='utf-8', errors='ignore')
    
    # 尝试解析JSON，格式化输出
    try:
        data = json.loads(json_content)
        formatted_text = json.dumps(data, ensure_ascii=False, indent=2)
    except json.JSONDecodeError:
        # 如果JSON格式错误，使用原始内容
        formatted_text = json_content
    
    return [
        Document(
            page_content=formatted_text,
            metadata={"source": path.name, "type": "json"}
        )
    ]

def load_image(path: Path) -> List[Document]:
    """
    加载图片文件,使用OCR提取文本内容
    使用Qwen2.5-VL多模态视觉模型与PaddleOCRVL模型进行图片理解,但是内存占用较高,直接卡死,后续待优化,暂时使用OCR模式识别
    """
    if not has_image_processing:
        print(f"⚠️ 跳过图片文件: {path.name} (未安装图片处理库)")
        return []
    
    docs = []
    
    try:
        # 使用PIL打开图片
        image = Image.open(path)
        
        # 预处理图片以提高OCR准确性
        # 1. 转换为灰度图
        image = image.convert('L')
        
        # 2. 调整对比度
        from PIL import ImageEnhance
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(4.0)  # 进一步增加对比度，提高数字识别
        
        # 3. 调整亮度
        enhancer = ImageEnhance.Brightness(image)
        image = enhancer.enhance(1.3)  # 稍微增加亮度
        
        # 4. 锐化处理，提高数字边缘清晰度
        enhancer = ImageEnhance.Sharpness(image)
        image = enhancer.enhance(2.0)  # 增加锐度
        
        # 5. 二值化处理，使用自适应阈值
        threshold = 140
        image = image.point(lambda x: 255 if x > threshold else 0, '1')
        
        # 6. 使用多种OCR配置尝试提取，特别优化数字识别
        configs = [
            '--psm 6 --oem 3 -c tessedit_char_whitelist=0123456789abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ：，。；',  # 统一块文本，指定字符集
            '--psm 3 --oem 3',  # 自动分段
            '--psm 11 --oem 3',  # 稀疏文本
        ]
        
        extracted_texts = []
        for config in configs:
            try:
                text = pytesseract.image_to_string(image, lang='chi_sim+eng', config=config)
                if text.strip():
                    extracted_texts.append(text)
            except Exception as e:
                print(f"⚠️ OCR配置 {config} 失败: {e}")
        
        # 合并所有提取的文本，去重和清理
        if extracted_texts:
            # 合并所有提取的文本
            combined_text = "\n".join(extracted_texts)
            
            # 清理文本
            cleaned_text = combined_text
            
            # 替换常见的OCR错误
            replacements = {
                ' . ': ' ',
                ' , ': ', ',
                ' : ': ': ',
                ' ; ': '; ',
                '  ': ' ',
                '。': '.',
                '，': ',',
                '：': ':',
                '；': ';',
                '　': ' ',  # 全角空格
                '贵用': '费用',
                '佐旅鼻': '差旅费',
                '业务招待鼻': '业务招待费',
            }
            
            for old, new in replacements.items():
                cleaned_text = cleaned_text.replace(old, new)
            
            # 去重
            lines = cleaned_text.strip().split('\n')
            unique_lines = []
            seen = set()
            for line in lines:
                line = line.strip()
                if line and line not in seen:
                    seen.add(line)
                    unique_lines.append(line)
            
            final_text = "\n".join(unique_lines)
            
            if final_text:
                print(f"✅ 从图片中提取文本: {final_text[:150]}...")
                docs.append(
                    Document(
                        page_content=final_text,
                        metadata={"source": path.name, "type": "image", "format": path.suffix.lower()[1:]}
                    )
                )
                print(f"✅ 从图片中提取文本成功: {path.name}")
                print(f"提取的完整文本: {final_text}")
            else:
                print(f"⚠️ 图片中未提取到文本: {path.name}")
        else:
            print(f"⚠️ 图片中未提取到文本: {path.name}")
            
    except Exception as e:
        print(f"⚠️ 处理图片文件时发生异常: {path.name}")
        print(f"💡 错误详情: {str(e)}")
    
    return docs


def ingest_files_0116(files: List[str], target="upload") -> int:
    
    docs: List[Document] = []
    for file in files:
        path = Path(file)
        if path.suffix.lower() == ".txt":
            docs.extend(load_txt(path))

        # 👉 后续你可以加 pdf / md / docx
        elif path.suffix.lower() == ".pdf":
            docs.extend(load_pdf(path))

    if not docs:
        return 0

    chunks = splitter.split_documents(docs)
    #VECTOR_DB.add_documents(chunks)
    print(f"📄 目标来源文件--------: {target}")
    if target == "kb":
        from rag_core import KB_VECTOR_DB
        KB_VECTOR_DB.add_documents(chunks)
    else:
        from rag_core import UPLOAD_VECTOR_DB
        UPLOAD_VECTOR_DB.add_documents(chunks)

    return len(chunks)

def ingest_files(
    paths: List[str],
    target: str = "kb",
    session_id: str = None
):
    """
    文件入库（支持 session / metadata / 训练沉淀）
    """
    all_chunks = []

    for path_str in paths:
        # 将字符串路径转换为Path对象
        path = Path(path_str)
        
        # 1️⃣ 加载文档
        #docs = load_files([path])
        ################################
        docs: List[Document] = []
        if path.suffix.lower() == ".txt":
            docs.extend(load_txt(path))
        elif path.suffix.lower() == ".md":
            # Markdown文件使用txt加载逻辑
            docs.extend(load_txt(path))
        elif path.suffix.lower() == ".pdf":
            docs.extend(load_pdf(path))
        elif path.suffix.lower() == ".docx":
            # Word文档处理
            if has_docx:
                doc = DocxDocument(path)
                text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
                docs.append(
                    Document(
                        page_content=text,
                        metadata={"source": path.name, "type": "docx"}
                    )
                )
            else:
                print(f"⚠️ 跳过.docx文件: {path.name} (未安装python-docx模块)")
        elif path.suffix.lower() == ".xlsx":
            # Excel文件处理
            docs.extend(load_xlsx(path))
        elif path.suffix.lower() == ".pptx":
            # PowerPoint文件处理
            docs.extend(load_pptx(path))
        elif path.suffix.lower() == ".csv":
            # CSV文件处理
            docs.extend(load_csv(path))
        elif path.suffix.lower() == ".html":
            # HTML文件处理
            docs.extend(load_html(path))
        elif path.suffix.lower() == ".json":
            # JSON文件处理
            docs.extend(load_json(path))
        elif path.suffix.lower() in [".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp"]:
            # 图片文件处理
            docs.extend(load_image(path))
        elif path.suffix.lower() == ".doc":
            # 旧版Word文档处理 - 多方法尝试处理
            try:
                # 方法1: 尝试使用python-docx2txt库（如果已安装）
                try:
                    import docx2txt
                    text = docx2txt.process(str(path))
                    if text.strip():
                        docs.append(
                            Document(
                                page_content=text,
                                metadata={"source": path.name, "type": "doc"}
                            )
                        )
                    else:
                        print(f"⚠️ 从.doc文件提取内容为空: {path.name}")
                except ImportError:
                    # 方法2: 尝试使用win32com（Windows系统可能可用）
                    try:
                        import win32com.client
                        word = win32com.client.Dispatch("Word.Application")
                        word.Visible = False
                        doc = word.Documents.Open(str(path))
                        text = doc.Content.Text
                        doc.Close()
                        word.Quit()
                        
                        if text.strip():
                            docs.append(
                                Document(
                                    page_content=text,
                                    metadata={"source": path.name, "type": "doc"}
                                )
                            )
                        else:
                            print(f"⚠️ 从.doc文件提取内容为空: {path.name}")
                    except Exception as e:
                        # 所有方法都失败，返回友好提示
                        print(f"⚠️ 无法处理.doc文件: {path.name}")
                        print("💡 提示: 请先将.doc文件转换为.docx或.txt格式后再上传")
                        print("💡 推荐: 可使用Microsoft Word或在线转换工具进行格式转换")
                        print(f"💡 详细错误: {str(e)}")
            except Exception as e:
                print(f"⚠️ 处理.doc文件时发生异常 : {path.name}")
                print(f"💡 错误详情: {str(e)}")
                print("💡 建议: 请先将.doc文件转换为.docx格式后再上传")
        else:
            print(f"⚠️ 跳过不支持的文件类型: {path.name}")
        ################################

        # 2️⃣ 文档级 metadata
        doc_id = uuid.uuid4().hex
        base_metadata = {
            "doc_id": doc_id,
            "source": target,               # kb / upload
            "session_id": session_id,       # None or session_id
            "filename": os.path.basename(path),
            "ingest_time": datetime.utcnow().isoformat()
        }
        
        # 调试信息
        print(f"   开始处理文件: {path.name}")
        print(f"   传入的session_id: {session_id}")
        print(f"   生成的base_metadata: {base_metadata}")
        print(f"   文档数量: {len(docs)}")

        for doc in docs:
            doc.metadata.update(base_metadata)
            # 调试信息
            print(f"   文档metadata更新后: {doc.metadata}")

        # 3️⃣ 切分
        #chunks = split_docs(docs)
        chunks = splitter.split_documents(docs)
        print(f"   切分后得到 {len(chunks)} 个chunk")

        # 4️⃣ chunk 级 metadata
        for idx, chunk in enumerate(chunks):
            chunk.metadata.update({
                **base_metadata,
                "chunk_id": f"{doc_id}_{idx}",
                "chunk_index": idx
            })
            print(f"   Chunk {idx} metadata: {chunk.metadata}")

        # 5️⃣ 写入向量库
        # vectorstore.add_documents(chunks)
        if chunks:  # 只在有chunks时才写入向量库
            print(f"   将 {len(chunks)} 个chunk添加到向量库")
            if target == "kb":
                from rag_core import KB_VECTOR_DB
                KB_VECTOR_DB.add_documents(chunks)
            else:
                from rag_core import UPLOAD_VECTOR_DB
                UPLOAD_VECTOR_DB.add_documents(chunks)
            print("   添加成功")
        else:
            print(f"⚠️ 没有从文件中提取到有效内容: {path.name}")

        all_chunks.extend(chunks)

    return {
        "files": len(paths),
        "chunks": len(all_chunks),
        "session_id": session_id,
        "target": target
    }






def ingest_files23(
    vectordb,
    files: List[Path]
) -> int:
    """
    接收文件路径列表，解析 + 切块 + 写入向量库
    """
    documents: List[Document] = []

    for path in files:
        if path.suffix.lower() == ".txt":
            documents.extend(load_txt(path))

        # 👉 后续你可以加 pdf / md / docx
        elif path.suffix.lower() == ".pdf":
            documents.extend(load_pdf(path))

    if not documents:
        return 0

    chunks = splitter.split_documents(documents)
    VECTOR_DB.add_documents(chunks)

    return len(chunks)


def ingest_old():
    all_docs = []

    for filename in os.listdir(DOCS_DIR):
        if filename.startswith("."):
            continue

        path = os.path.join(DOCS_DIR, filename)
        ext = filename.lower().split(".")[-1]

        if ext == "txt":
            all_docs.extend(load_txt_old(path))

        elif ext == "md":
            all_docs.extend(load_md(path))

        elif ext == "pdf":
            all_docs.extend(load_pdf(path))

        elif ext == "docx":
            all_docs.extend(load_docx(path))

        else:
            print(f"⚠️ 跳过不支持的文件类型: {filename}")

    if not all_docs:
        print("⚠️ 没有加载到任何文档")
        return

    chunks = splitter.split_documents(all_docs)
    VECTOR_DB.add_documents(chunks)
    #vectordb.persist()

    print(f"✅ 新增入库 chunk 数量: {len(chunks)}")

#if __name__ == "__main__":
#    ingest_old()
